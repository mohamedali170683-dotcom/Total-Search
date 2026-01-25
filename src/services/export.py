"""Export service for generating PDF and PowerPoint reports."""

import io
from datetime import datetime
from typing import Any

from reportlab.lib import colors
from reportlab.lib.pagesizes import letter, A4
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.platypus import (
    SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, Image, PageBreak
)
from reportlab.graphics.shapes import Drawing, Rect
from reportlab.graphics.charts.piecharts import Pie
from reportlab.graphics.charts.barcharts import VerticalBarChart

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RgbColor
    from pptx.enum.text import PP_ALIGN
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


PLATFORM_COLORS = {
    "google": colors.HexColor("#4285F4"),
    "youtube": colors.HexColor("#FF0000"),
    "amazon": colors.HexColor("#FF9900"),
    "tiktok": colors.HexColor("#000000"),
    "instagram": colors.HexColor("#E1306C"),
    "pinterest": colors.HexColor("#E60023"),
}

PLATFORM_COLORS_HEX = {
    "google": "#4285F4",
    "youtube": "#FF0000",
    "amazon": "#FF9900",
    "tiktok": "#000000",
    "instagram": "#E1306C",
    "pinterest": "#E60023",
}


class DemandReportPDF:
    """Generate PDF reports for demand analysis."""

    def __init__(self):
        self.styles = getSampleStyleSheet()
        self._add_custom_styles()

    def _add_custom_styles(self):
        """Add custom paragraph styles."""
        self.styles.add(ParagraphStyle(
            name="ReportTitle",
            parent=self.styles["Heading1"],
            fontSize=24,
            spaceAfter=30,
            textColor=colors.HexColor("#1a1a1a"),
        ))
        self.styles.add(ParagraphStyle(
            name="SectionTitle",
            parent=self.styles["Heading2"],
            fontSize=16,
            spaceBefore=20,
            spaceAfter=12,
            textColor=colors.HexColor("#333333"),
        ))
        self.styles.add(ParagraphStyle(
            name="Insight",
            parent=self.styles["Normal"],
            fontSize=11,
            leftIndent=20,
            spaceBefore=6,
            spaceAfter=6,
            textColor=colors.HexColor("#444444"),
        ))

    def generate(
        self,
        analysis_data: dict,
        title: str = "Demand Distribution Analysis",
    ) -> bytes:
        """
        Generate a PDF report from demand analysis data.

        Args:
            analysis_data: Data from /api/demand/analyze endpoint
            title: Report title

        Returns:
            PDF bytes
        """
        buffer = io.BytesIO()
        doc = SimpleDocTemplate(
            buffer,
            pagesize=letter,
            rightMargin=72,
            leftMargin=72,
            topMargin=72,
            bottomMargin=72,
        )

        elements = []

        # Title
        elements.append(Paragraph(title, self.styles["ReportTitle"]))
        elements.append(Paragraph(
            f"Generated: {datetime.utcnow().strftime('%B %d, %Y')}",
            self.styles["Normal"]
        ))
        elements.append(Spacer(1, 20))

        # Executive Summary
        elements.append(Paragraph("Executive Summary", self.styles["SectionTitle"]))
        summary = analysis_data.get("distribution_summary", {})
        keywords = analysis_data.get("keywords", [])

        summary_text = f"""
        Analysis of {len(keywords)} keyword(s): {', '.join(keywords[:5])}{'...' if len(keywords) > 5 else ''}
        <br/><br/>
        Total Demand Volume: <b>{analysis_data.get('total_demand', 0):,}</b>
        <br/>
        Google Share: <b>{summary.get('google_share', 0)}%</b> |
        Non-Google Share: <b>{summary.get('non_google_share', 0)}%</b>
        <br/>
        Top Platform: <b>{summary.get('top_platform', 'N/A').title()}</b>
        """
        elements.append(Paragraph(summary_text, self.styles["Normal"]))
        elements.append(Spacer(1, 20))

        # Platform Breakdown Table
        elements.append(Paragraph("Platform Breakdown", self.styles["SectionTitle"]))
        platform_data = analysis_data.get("platforms", [])

        if platform_data:
            table_data = [["Platform", "Volume", "Share", "Trend"]]
            for p in platform_data:
                table_data.append([
                    p.get("display_name", p.get("platform", "Unknown")),
                    f"{p.get('volume', 0):,}",
                    f"{p.get('percentage', 0)}%",
                    p.get("trend", "stable").title(),
                ])

            table = Table(table_data, colWidths=[2*inch, 1.5*inch, 1*inch, 1*inch])
            table.setStyle(TableStyle([
                ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#4285F4")),
                ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
                ("ALIGN", (0, 0), (-1, -1), "CENTER"),
                ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                ("FONTSIZE", (0, 0), (-1, 0), 11),
                ("BOTTOMPADDING", (0, 0), (-1, 0), 12),
                ("BACKGROUND", (0, 1), (-1, -1), colors.HexColor("#f8f9fa")),
                ("GRID", (0, 0), (-1, -1), 1, colors.HexColor("#dee2e6")),
                ("FONTSIZE", (0, 1), (-1, -1), 10),
                ("TOPPADDING", (0, 1), (-1, -1), 8),
                ("BOTTOMPADDING", (0, 1), (-1, -1), 8),
            ]))
            elements.append(table)
            elements.append(Spacer(1, 20))

        # Pie Chart
        if platform_data:
            elements.append(Paragraph("Demand Distribution", self.styles["SectionTitle"]))
            drawing = self._create_pie_chart(platform_data)
            elements.append(drawing)
            elements.append(Spacer(1, 20))

        # Insights
        insights = analysis_data.get("insights", [])
        if insights:
            elements.append(Paragraph("Key Insights", self.styles["SectionTitle"]))
            for insight in insights:
                icon = "●" if insight.get("priority") == "high" else "○"
                elements.append(Paragraph(
                    f"{icon} <b>{insight.get('title', '')}</b>",
                    self.styles["Insight"]
                ))
                elements.append(Paragraph(
                    insight.get("description", ""),
                    self.styles["Insight"]
                ))
            elements.append(Spacer(1, 20))

        # Recommendations
        recommendations = analysis_data.get("recommendations", [])
        if recommendations:
            elements.append(Paragraph("Recommendations", self.styles["SectionTitle"]))
            for i, rec in enumerate(recommendations, 1):
                elements.append(Paragraph(
                    f"<b>{i}. {rec.get('action', '')}</b>",
                    self.styles["Insight"]
                ))
                elements.append(Paragraph(
                    rec.get("description", ""),
                    self.styles["Insight"]
                ))
                tactics = rec.get("tactics", [])
                if tactics:
                    tactics_text = "Tactics: " + ", ".join(tactics)
                    elements.append(Paragraph(tactics_text, self.styles["Insight"]))
                elements.append(Spacer(1, 6))

        # Footer
        elements.append(Spacer(1, 30))
        elements.append(Paragraph(
            "Report generated by Total Search - Cross-Platform Demand Intelligence",
            self.styles["Normal"]
        ))

        doc.build(elements)
        buffer.seek(0)
        return buffer.getvalue()

    def _create_pie_chart(self, platform_data: list, width: int = 400, height: int = 200) -> Drawing:
        """Create a pie chart for platform distribution."""
        drawing = Drawing(width, height)

        pie = Pie()
        pie.x = 100
        pie.y = 20
        pie.width = 150
        pie.height = 150

        pie.data = [p.get("volume", 0) for p in platform_data if p.get("volume", 0) > 0]
        pie.labels = [p.get("display_name", p.get("platform", "")) for p in platform_data if p.get("volume", 0) > 0]

        # Set colors
        for i, p in enumerate(platform_data):
            if p.get("volume", 0) > 0:
                platform = p.get("platform", "google")
                pie.slices[i].fillColor = PLATFORM_COLORS.get(platform, colors.gray)

        pie.slices.strokeWidth = 1
        pie.slices.strokeColor = colors.white

        drawing.add(pie)
        return drawing


class DemandReportPPTX:
    """Generate PowerPoint reports for demand analysis."""

    def __init__(self):
        if not PPTX_AVAILABLE:
            raise ImportError("python-pptx is required for PowerPoint export")

    def generate(
        self,
        analysis_data: dict,
        title: str = "Demand Distribution Analysis",
    ) -> bytes:
        """
        Generate a PowerPoint report from demand analysis data.

        Args:
            analysis_data: Data from /api/demand/analyze endpoint
            title: Report title

        Returns:
            PPTX bytes
        """
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # Title Slide
        self._add_title_slide(prs, title, analysis_data)

        # Executive Summary Slide
        self._add_summary_slide(prs, analysis_data)

        # Platform Breakdown Slide
        self._add_platform_slide(prs, analysis_data)

        # Insights Slide
        self._add_insights_slide(prs, analysis_data)

        # Recommendations Slide
        self._add_recommendations_slide(prs, analysis_data)

        buffer = io.BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        return buffer.getvalue()

    def _add_title_slide(self, prs: "Presentation", title: str, data: dict):
        """Add title slide."""
        slide_layout = prs.slide_layouts[6]  # Blank
        slide = prs.slides.add_slide(slide_layout)

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12), Inches(1.5))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(44)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        # Subtitle
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(12), Inches(0.5))
        tf = subtitle_box.text_frame
        p = tf.paragraphs[0]
        keywords = data.get("keywords", [])
        p.text = f"Analysis of: {', '.join(keywords[:5])}{'...' if len(keywords) > 5 else ''}"
        p.font.size = Pt(20)
        p.alignment = PP_ALIGN.CENTER

        # Date
        date_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12), Inches(0.5))
        tf = date_box.text_frame
        p = tf.paragraphs[0]
        p.text = datetime.utcnow().strftime("%B %d, %Y")
        p.font.size = Pt(14)
        p.alignment = PP_ALIGN.CENTER

    def _add_summary_slide(self, prs: "Presentation", data: dict):
        """Add executive summary slide."""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "Executive Summary"
        p.font.size = Pt(32)
        p.font.bold = True

        summary = data.get("distribution_summary", {})

        # Key metrics
        metrics = [
            ("Total Demand", f"{data.get('total_demand', 0):,}"),
            ("Google Share", f"{summary.get('google_share', 0)}%"),
            ("Non-Google Share", f"{summary.get('non_google_share', 0)}%"),
            ("Top Platform", summary.get('top_platform', 'N/A').title()),
        ]

        for i, (label, value) in enumerate(metrics):
            x = Inches(0.5 + (i * 3.2))
            box = slide.shapes.add_textbox(x, Inches(1.5), Inches(3), Inches(1.5))
            tf = box.text_frame

            p = tf.paragraphs[0]
            p.text = value
            p.font.size = Pt(36)
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER

            p2 = tf.add_paragraph()
            p2.text = label
            p2.font.size = Pt(14)
            p2.alignment = PP_ALIGN.CENTER

    def _add_platform_slide(self, prs: "Presentation", data: dict):
        """Add platform breakdown slide."""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "Platform Breakdown"
        p.font.size = Pt(32)
        p.font.bold = True

        platforms = data.get("platforms", [])

        # Create table
        if platforms:
            rows = len(platforms) + 1
            cols = 4
            table = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.2), Inches(12), Inches(0.5 * rows)).table

            # Headers
            headers = ["Platform", "Volume", "Share", "Trend"]
            for i, header in enumerate(headers):
                cell = table.cell(0, i)
                cell.text = header
                cell.text_frame.paragraphs[0].font.bold = True
                cell.text_frame.paragraphs[0].font.size = Pt(14)

            # Data
            for row_idx, p in enumerate(platforms, 1):
                table.cell(row_idx, 0).text = p.get("display_name", p.get("platform", ""))
                table.cell(row_idx, 1).text = f"{p.get('volume', 0):,}"
                table.cell(row_idx, 2).text = f"{p.get('percentage', 0)}%"
                table.cell(row_idx, 3).text = p.get("trend", "stable").title()

    def _add_insights_slide(self, prs: "Presentation", data: dict):
        """Add insights slide."""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "Key Insights"
        p.font.size = Pt(32)
        p.font.bold = True

        insights = data.get("insights", [])

        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12), Inches(5.5))
        tf = content_box.text_frame

        for i, insight in enumerate(insights[:5]):
            if i > 0:
                p = tf.add_paragraph()
            else:
                p = tf.paragraphs[0]

            p.text = f"• {insight.get('title', '')}"
            p.font.size = Pt(18)
            p.font.bold = True
            p.space_after = Pt(6)

            p2 = tf.add_paragraph()
            p2.text = f"  {insight.get('description', '')}"
            p2.font.size = Pt(14)
            p2.space_after = Pt(18)

    def _add_recommendations_slide(self, prs: "Presentation", data: dict):
        """Add recommendations slide."""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "Recommendations"
        p.font.size = Pt(32)
        p.font.bold = True

        recommendations = data.get("recommendations", [])

        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12), Inches(5.5))
        tf = content_box.text_frame

        for i, rec in enumerate(recommendations[:4]):
            if i > 0:
                p = tf.add_paragraph()
            else:
                p = tf.paragraphs[0]

            p.text = f"{i+1}. {rec.get('action', '')}"
            p.font.size = Pt(20)
            p.font.bold = True
            p.space_after = Pt(6)

            p2 = tf.add_paragraph()
            p2.text = f"   {rec.get('description', '')}"
            p2.font.size = Pt(14)
            p2.space_after = Pt(12)


def generate_demand_report_pdf(analysis_data: dict, title: str = "Demand Distribution Analysis") -> bytes:
    """Generate a PDF demand report."""
    generator = DemandReportPDF()
    return generator.generate(analysis_data, title)


def generate_demand_report_pptx(analysis_data: dict, title: str = "Demand Distribution Analysis") -> bytes:
    """Generate a PowerPoint demand report."""
    generator = DemandReportPPTX()
    return generator.generate(analysis_data, title)
